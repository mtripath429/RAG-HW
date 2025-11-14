# parsers.py
import os
from typing import Optional

import pandas as pd
from pypdf import PdfReader
import docx
from pptx import Presentation


def parse_pdf(path: str) -> str:
    reader = PdfReader(path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)


def parse_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def parse_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_md(path: str) -> str:
    # Treat markdown as plain text
    return parse_txt(path)


def parse_csv(path: str) -> str:
    df = pd.read_csv(path)
    # Serialize the table to CSV text for RAG
    return df.to_csv(index=False)


def parse_file(path: str) -> Optional[str]:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return parse_pdf(path)
        elif ext == ".docx":
            return parse_docx(path)
        elif ext == ".pptx":
            return parse_pptx(path)
        elif ext == ".txt":
            return parse_txt(path)
        elif ext in [".md", ".markdown"]:
            return parse_md(path)
        elif ext == ".csv":
            return parse_csv(path)
        else:
            # Unsupported file type for now
            return None
    except Exception as e:
        print(f"Error parsing {path}: {e}")
        return None
